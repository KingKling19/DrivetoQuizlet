#!/usr/bin/env python3
"""
convert_folder_to_quizlet.py

Usage:
  # Process all .pptx files in a folder
  python convert_folder_to_quizlet.py "C:/path/to/slides"

  # Or process a single file
  python convert_folder_to_quizlet.py "C:/path/to/slides/mydeck.pptx"

Options:
  --model gpt-4o-mini          Model to use (default: gpt-4o-mini)
  --window 3                   Slides per LLM window (default: 3)
  --max-retries 4              API retry attempts (default: 4)
  --sleep 0.3                  Base sleep between API calls in seconds (default: 0.3)
  --min-def-len 12             Drop defs shorter than this (default: 12)
  --dry-run                    Parse slides but do not call the API
  --budget-usd 0               Optional soft budget (USD). If >0, exit when estimate exceeded.
  --verbose                    Extra logs
Env:
  OPENAI_API_KEY must be set.
"""

import argparse
import csv
import json
import os
import re
import sys
import time
from collections import defaultdict
from pathlib import Path
from typing import List, Dict, Any, Optional
import pickle
import numpy as np

try:
    from openai import OpenAI
except Exception as e:
    print("ERROR: openai>=1.0.0 is required. pip install openai", file=sys.stderr)
    raise

try:
    from pptx import Presentation
except Exception as e:
    print("ERROR: python-pptx is required. pip install python-pptx", file=sys.stderr)
    raise


# ---------------------------
# Configuration defaults
# ---------------------------
DEFAULT_MODEL = "gpt-4o-mini"
SYSTEM_PROMPT = """You extract high-quality, testable term–definition pairs or big-idea summaries from training slides.
Return ONLY valid JSON (list of objects). No commentary.

Guidelines:
- Prioritize the main concepts, big ideas, and key takeaways from each slide over isolated words.
- Include acronyms ONLY if they are explicitly defined or expanded in the provided text; otherwise, assume acronyms are already understood and do not output separate cards for them.
- Prefer doctrinal phrases, bold/ALL-CAPS titles, and glossary-like lines when they convey an important concept.
- The definition MUST be fully supported by the provided text (do not invent).
- Keep definitions concise (1–2 sentences).
- Skip generic or trivial bullets (e.g., "Objectives", "Agenda", "Summary").
- If the context is ambiguous, omit the item rather than guessing.
- Use cross-lesson context to enhance understanding and avoid duplicate definitions.
- Consider related concepts from other lessons when creating definitions.

Output schema:
[{"term": str, "definition": str, "source_slide": int, "confidence": float}]
"""

SKIP_TITLES = {"objectives", "agenda", "summary", "references", "q&a", "questions"}
MIN_DEF_LEN = 12

# Cross-lesson context configuration
CONTEXT_CONFIG = {
    "max_related_lessons": 3,
    "context_weight_threshold": 0.3,
    "max_context_length": 2000,
    "include_prerequisites": True,
    "include_related_concepts": True
}


# ---------------------------
# Slide parsing
# ---------------------------
def _clean_ws(s: str) -> str:
    s = re.sub(r"\r\n?", "\n", s or "")
    s = re.sub(r"[ \t]+", " ", s)
    s = re.sub(r"\n{2,}", "\n", s)
    return s.strip()


def extract_slide_text(slide) -> Dict[str, str]:
    # Title
    title = ""
    try:
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            title = slide.shapes.title.text or ""
    except Exception:
        pass

    # Bodies & tables
    bodies = []
    for shape in slide.shapes:
        try:
            if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                t = shape.text or ""
                if t and t != title:
                    bodies.append(t)
            if getattr(shape, "has_table", False) and shape.has_table:
                tbl_text = []
                for r in shape.table.rows:
                    cells = []
                    for c in r.cells:
                        cells.append((c.text or "").strip())
                    tbl_text.append(" | ".join(cells))
                if tbl_text:
                    bodies.append("\n".join(tbl_text))
        except Exception:
            # skip shapes we can't read
            continue

    # Speaker notes
    notes = ""
    try:
        if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text or ""
    except Exception:
        pass

    return {
        "title": _clean_ws(title),
        "body": _clean_ws("\n".join(bodies)),
        "notes": _clean_ws(notes),
    }


def slide_windows(slides_ctx: List[Dict[str, Any]], window: int = 3) -> List[List[Dict[str, Any]]]:
    """
    Overlapping windows to let the model see neighboring slides where defs often live.
    Example with window=3: [0..2], [2..4], [4..6] (stride = window-1)
    """
    i = 0
    chunks = []
    stride = max(1, window - 1)
    while i < len(slides_ctx):
        chunk = slides_ctx[i:i + window]
        if chunk:
            chunks.append(chunk)
        i += stride
    return chunks


def build_prompt(chunk: List[Dict[str, Any]], cross_lesson_context: str = "") -> str:
    parts = []
    for s in chunk:
        title_line = f"# Slide {s['index']}: {s['title']}".strip()
        parts.append(title_line)
        if s["body"]:
            parts.append(s["body"])
        if s["notes"]:
            parts.append(f"(Notes) {s['notes']}")
        parts.append("")  # spacer
    ctx = "\n".join(parts).strip()
    
    prompt = "Extract term–definition pairs from the following slides. Return ONLY JSON (no markdown)."
    
    if cross_lesson_context:
        prompt += f"\n\nCross-Lesson Context:\n{cross_lesson_context}\n"
    
    prompt += f"\n\nSlides:\n{ctx}"
    return prompt


def load_cross_lesson_data(config_dir: Path = Path("config")) -> Dict[str, Any]:
    """Load cross-lesson analysis data for context enhancement."""
    data = {
        "content_index": {},
        "semantic_embeddings": {},
        "lesson_relationships": {},
        "cross_references": {}
    }
    
    try:
        # Load content index
        index_file = config_dir / "lesson_content_index.json"
        if index_file.exists():
            with open(index_file, 'r', encoding='utf-8') as f:
                data["content_index"] = json.load(f)
        
        # Load semantic embeddings
        embeddings_file = config_dir / "semantic_embeddings.pkl"
        if embeddings_file.exists():
            with open(embeddings_file, 'rb') as f:
                data["semantic_embeddings"] = pickle.load(f)
        
        # Load lesson relationships
        relationships_file = config_dir / "lesson_relationships_analysis.json"
        if relationships_file.exists():
            with open(relationships_file, 'r', encoding='utf-8') as f:
                data["lesson_relationships"] = json.load(f)
        
        # Load cross-references
        cross_refs_file = config_dir / "cross_references.json"
        if cross_refs_file.exists():
            with open(cross_refs_file, 'r', encoding='utf-8') as f:
                data["cross_references"] = json.load(f)
        
        print(f"✓ Loaded cross-lesson data: {len(data['content_index'])} lessons indexed")
        return data
    except Exception as e:
        print(f"⚠️  Could not load cross-lesson data: {e}")
        return data


def get_lesson_id_from_path(pptx_path: Path) -> str:
    """Extract lesson ID from presentation path."""
    # Try to find lesson directory in path
    path_parts = pptx_path.parts
    for i, part in enumerate(path_parts):
        if part == "lessons" and i + 1 < len(path_parts):
            return path_parts[i + 1]
    
    # Fallback: use filename without extension
    return pptx_path.stem


def find_related_lessons(lesson_id: str, cross_lesson_data: Dict[str, Any], max_lessons: int = 3) -> List[Dict[str, Any]]:
    """Find lessons related to the current lesson for context enhancement."""
    related_lessons = []
    
    try:
        relationships = cross_lesson_data.get("lesson_relationships", {})
        lesson_rels = relationships.get(lesson_id, {})
        
        # Get related lessons - handle both old and new formats
        related = lesson_rels.get("related_lessons", [])
        
        # If related_lessons is a list of strings (old format), convert to new format
        if related and isinstance(related[0], str):
            # Old format - convert to new format
            related_lessons_new = []
            for rel_id in related:
                similarity = lesson_rels.get("relationship_scores", {}).get(rel_id, 0.0)
                related_lessons_new.append({
                    "lesson_id": rel_id,
                    "similarity_score": similarity,
                    "relationship_type": "related",
                    "related_concepts": []
                })
            related = related_lessons_new
        
        # Sort by similarity score
        if related and isinstance(related[0], dict):
            related.sort(key=lambda x: x.get("similarity_score", 0), reverse=True)
        
        # Take top related lessons
        for rel in related[:max_lessons]:
            if isinstance(rel, dict):
                related_lesson_id = rel.get("lesson_id")
                if related_lesson_id and related_lesson_id != lesson_id:
                    related_lessons.append(rel)
        
        return related_lessons
    except Exception as e:
        print(f"⚠️  Error finding related lessons: {e}")
        return []


def extract_related_context(related_lessons: List[Dict[str, Any]], cross_lesson_data: Dict[str, Any]) -> str:
    """Extract relevant context from related lessons."""
    context_parts = []
    
    try:
        content_index = cross_lesson_data.get("content_index", {})
        
        for rel in related_lessons:
            if not isinstance(rel, dict):
                continue
                
            lesson_id = rel.get("lesson_id")
            similarity_score = rel.get("similarity_score", 0)
            related_concepts = rel.get("related_concepts", [])
            
            if lesson_id in content_index:
                lesson_data = content_index[lesson_id]
                if not isinstance(lesson_data, dict):
                    continue
                    
                lesson_name = lesson_data.get("lesson_name", lesson_id)
                
                # Add lesson header
                context_parts.append(f"## Related Lesson: {lesson_name} (Similarity: {similarity_score:.2f})")
                
                # Add key concepts
                if related_concepts and isinstance(related_concepts, list):
                    context_parts.append("### Key Related Concepts:")
                    for concept in related_concepts[:5]:  # Limit to top 5 concepts
                        context_parts.append(f"- {concept}")
                
                # Add content snippets from presentations
                content_sources = lesson_data.get("content_sources", {})
                if isinstance(content_sources, dict):
                    presentations = content_sources.get("presentations", {})
                    if isinstance(presentations, dict):
                        for pptx_name, pptx_data in list(presentations.items())[:2]:  # Limit to 2 presentations
                            if isinstance(pptx_data, dict):
                                slides = pptx_data.get("slides", [])
                                if isinstance(slides, list):
                                    for slide in slides[:3]:  # Limit to 3 slides per presentation
                                        if isinstance(slide, dict):
                                            title = slide.get("title", "")
                                            body = slide.get("body", "")
                                            if title and body:
                                                context_parts.append(f"### {title}")
                                                context_parts.append(body[:300] + "..." if len(body) > 300 else body)
                
                context_parts.append("")  # Spacer
        
        context = "\n".join(context_parts)
        
        # Limit total context length
        if len(context) > CONTEXT_CONFIG["max_context_length"]:
            context = context[:CONTEXT_CONFIG["max_context_length"]] + "..."
        
        return context
    except Exception as e:
        print(f"⚠️  Error extracting related context: {e}")
        return ""


def enhance_with_cross_lesson_context(chunk: List[Dict[str, Any]], pptx_path: Path, cross_lesson_data: Dict[str, Any]) -> str:
    """Enhance slide content with cross-lesson context."""
    try:
        lesson_id = get_lesson_id_from_path(pptx_path)
        related_lessons = find_related_lessons(lesson_id, cross_lesson_data, CONTEXT_CONFIG["max_related_lessons"])
        
        if not related_lessons:
            return ""
        
        return extract_related_context(related_lessons, cross_lesson_data)
    except Exception as e:
        print(f"⚠️  Error enhancing with cross-lesson context: {e}")
        return ""


# ---------------------------
# LLM calling & post-processing
# ---------------------------
def call_llm(client: OpenAI, model: str, prompt: str, max_retries: int = 4, base_sleep: float = 0.3, verbose: bool = False):
    last_err = None
    for attempt in range(max_retries):
        try:
            resp = client.chat.completions.create(
                model=model,
                temperature=0.15,
                messages=[
                    {"role": "system", "content": SYSTEM_PROMPT},
                    {"role": "user", "content": prompt},
                ],
            )
            txt = (resp.choices[0].message.content or "").strip()
            # Strip codefences if any
            txt = re.sub(r"^```(?:json)?|```$", "", txt, flags=re.MULTILINE).strip()
            data = json.loads(txt)
            if isinstance(data, list):
                return data
            # If not a list, try to coerce
            return []
        except json.JSONDecodeError as e:
            last_err = e
            if verbose:
                print(f"[warn] JSON decode failed (attempt {attempt+1}/{max_retries}).", file=sys.stderr)
        except Exception as e:
            last_err = e
            msg = str(e).lower()
            if any(x in msg for x in ("rate", "quota", "429", "timeout", "temporarily unavailable", "overloaded", "500", "502", "503", "504")):
                # exponential backoff
                sleep_s = base_sleep * (2 ** attempt)
                if verbose:
                    print(f"[info] transient error: {e}. sleeping {sleep_s:.2f}s", file=sys.stderr)
                time.sleep(sleep_s)
                continue
            if verbose:
                print(f"[error] non-retryable error: {e}", file=sys.stderr)
            break
    if verbose and last_err:
        print(f"[error] giving up after {max_retries} tries: {last_err}", file=sys.stderr)
    return []


def is_skippable_title(title: str) -> bool:
    if not title:
        return False
    t = re.sub(r"[^a-z]", "", title.lower())
    return t in SKIP_TITLES


def good_definition(defi: str, min_len: int) -> bool:
    d = (defi or "").strip()
    if len(d) < min_len:
        return False
    if len(d) > 600:
        return False
    return True


def canonical_term(t: str) -> str:
    """Enhanced term normalization with better handling of variations."""
    t = (t or "").strip()
    # Remove common prefixes/suffixes that don't change meaning
    t = re.sub(r"^(the|a|an)\s+", "", t.lower())
    t = re.sub(r"\s+(system|process|procedure|method|technique)$", "", t)
    # Normalize whitespace and punctuation
    t = re.sub(r"\s+", " ", t)
    t = re.sub(r"[^\w\s]", "", t)  # Remove punctuation
    return t.strip()


def compact_definition(d: str) -> str:
    """Enhanced definition compression with better sentence selection."""
    # keep 1–2 sentences, prioritizing complete thoughts
    parts = re.split(r"(?<=[.!?])\s+", d.strip())
    if len(parts) <= 2:
        return " ".join(parts).strip()
    
    # Select sentences with more content (longer sentences)
    scored_parts = [(part, len(part.split())) for part in parts if part.strip()]
    scored_parts.sort(key=lambda x: x[1], reverse=True)
    
    # Take the two most substantial sentences
    selected = [part for part, _ in scored_parts[:2]]
    selected.sort(key=lambda x: parts.index(x))  # Maintain original order
    
    return " ".join(selected).strip()


def calculate_semantic_similarity(text1: str, text2: str) -> float:
    """Calculate semantic similarity between two texts using simple word overlap."""
    if not text1 or not text2:
        return 0.0
    
    # Normalize texts
    text1_words = set(re.findall(r'\b\w+\b', text1.lower()))
    text2_words = set(re.findall(r'\b\w+\b', text2.lower()))
    
    if not text1_words or not text2_words:
        return 0.0
    
    # Calculate Jaccard similarity
    intersection = len(text1_words & text2_words)
    union = len(text1_words | text2_words)
    
    return intersection / union if union > 0 else 0.0


def calculate_fuzzy_similarity(term1: str, term2: str) -> float:
    """Calculate fuzzy similarity between terms using Levenshtein distance."""
    if not term1 or not term2:
        return 0.0
    
    # Simple implementation - can be enhanced with proper Levenshtein
    term1_norm = canonical_term(term1)
    term2_norm = canonical_term(term2)
    
    if term1_norm == term2_norm:
        return 1.0
    
    # Calculate character-level similarity
    max_len = max(len(term1_norm), len(term2_norm))
    if max_len == 0:
        return 0.0
    
    # Simple character overlap
    common_chars = sum(1 for c in term1_norm if c in term2_norm)
    return common_chars / max_len


def is_duplicate_flashcard(card1: Dict[str, Any], card2: Dict[str, Any], config: Dict[str, Any]) -> bool:
    """Enhanced duplicate detection using multiple similarity metrics."""
    term1 = card1.get("term", "").strip()
    term2 = card2.get("term", "").strip()
    def1 = card1.get("definition", "").strip()
    def2 = card2.get("definition", "").strip()
    
    if not term1 or not term2:
        return False
    
    # 1. Exact term match
    if canonical_term(term1) == canonical_term(term2):
        return True
    
    # 2. Fuzzy term similarity
    fuzzy_threshold = config.get("fuzzy_match_threshold", 0.3)
    if calculate_fuzzy_similarity(term1, term2) > fuzzy_threshold:
        return True
    
    # 3. Semantic similarity (if terms are similar and definitions are very similar)
    semantic_threshold = config.get("semantic_similarity_threshold", 0.85)
    term_similarity = calculate_fuzzy_similarity(term1, term2)
    def_similarity = calculate_semantic_similarity(def1, def2)
    
    if term_similarity > 0.5 and def_similarity > semantic_threshold:
        return True
    
    return False


def select_best_duplicate(card1: Dict[str, Any], card2: Dict[str, Any], config: Dict[str, Any]) -> Dict[str, Any]:
    """Select the better flashcard when duplicates are found."""
    conf1 = float(card1.get("confidence", 0.0) or 0.0)
    conf2 = float(card2.get("confidence", 0.0) or 0.0)
    
    # Compare confidence scores
    if abs(conf1 - conf2) > 0.1:
        return card1 if conf1 > conf2 else card2
    
    # Compare definition quality (length and structure)
    def1 = card1.get("definition", "")
    def2 = card2.get("definition", "")
    
    # Prefer longer, more detailed definitions
    if len(def1) > len(def2) * 1.2:
        return card1
    elif len(def2) > len(def1) * 1.2:
        return card2
    
    # Prefer definitions with better structure
    def1_score = len([s for s in re.split(r'[.!?]+', def1) if len(s.strip().split()) >= 3])
    def2_score = len([s for s in re.split(r'[.!?]+', def2) if len(s.strip().split()) >= 3])
    
    if def1_score != def2_score:
        return card1 if def1_score > def2_score else card2
    
    # Default to first card if all else is equal
    return card1


def dedupe_and_filter(cards: List[Dict[str, Any]], min_def_len: int, config: Optional[Dict[str, Any]] = None) -> List[Dict[str, Any]]:
    """Enhanced duplicate detection and filtering with semantic similarity."""
    if config is None:
        config = {
            "fuzzy_match_threshold": 0.3,
            "semantic_similarity_threshold": 0.85,
            "context_weight": 0.2,
            "confidence_weight": 0.8
        }
    
    # Filter out invalid cards first
    valid_cards = []
    for c in cards:
        term = (c.get("term") or "").strip()
        defi = (c.get("definition") or "").strip()
        
        if term and good_definition(defi, min_def_len):
            valid_cards.append(c)
    
    if not valid_cards:
        return []
    
    # Enhanced duplicate detection
    unique_cards = []
    processed_terms = set()
    
    for i, card in enumerate(valid_cards):
        term = card.get("term", "").strip()
        canonical_key = canonical_term(term)
        
        # Check if this is a duplicate of any existing card
        is_duplicate = False
        best_card = card
        
        for existing_card in unique_cards:
            if is_duplicate_flashcard(card, existing_card, config):
                is_duplicate = True
                # Select the better card
                best_card = select_best_duplicate(card, existing_card, config)
                
                # Remove the existing card if it's not the best
                if best_card != existing_card:
                    unique_cards.remove(existing_card)
                    unique_cards.append(best_card)
                break
        
        if not is_duplicate:
            unique_cards.append(card)
    
    # Final processing
    result = []
    for card in unique_cards:
        term = card.get("term", "").strip()
        defi = card.get("definition", "").strip()
        conf = float(card.get("confidence", 0.0) or 0.0)
        src = int(card.get("source_slide", 0) or 0)
        
        result.append({
            "term": term,
            "definition": compact_definition(defi),
            "confidence": conf,
            "source_slide": src,
        })
    
    # Sort: source_slide then term
    return sorted(result, key=lambda x: (x["source_slide"], x["term"].lower()))


def extract_cards_from_presentation(pptx_path: Path, client: OpenAI, model: str, window: int, max_retries: int, base_sleep: float, min_def_len: int, dry_run: bool, verbose: bool, cross_lesson_data: Optional[Dict[str, Any]] = None) -> List[Dict[str, Any]]:
    prs = Presentation(str(pptx_path))
    slides_ctx = []
    for idx, slide in enumerate(prs.slides, start=1):
        info = extract_slide_text(slide)
        info["index"] = idx
        slides_ctx.append(info)

    # Optionally filter out housekeeping-only slides
    filtered = []
    for s in slides_ctx:
        if is_skippable_title(s["title"]) and not s["body"] and not s["notes"]:
            continue
        filtered.append(s)

    chunks = slide_windows(filtered, window=window)
    all_cards: List[Dict[str, Any]] = []

    for chunk in chunks:
        # Enhance with cross-lesson context if available
        cross_lesson_context = ""
        if cross_lesson_data:
            cross_lesson_context = enhance_with_cross_lesson_context(chunk, pptx_path, cross_lesson_data)
            if verbose and cross_lesson_context:
                print(f"[context] Enhanced chunk with cross-lesson context ({len(cross_lesson_context)} chars)")
        
        prompt = build_prompt(chunk, cross_lesson_context)
        if dry_run:
            if verbose:
                print(f"[dry-run] would call LLM on slides {[c['index'] for c in chunk]}")
            continue

        data = call_llm(client, model, prompt, max_retries=max_retries, base_sleep=base_sleep, verbose=verbose)
        if not data:
            continue

        # Fill missing source_slide with the chunk's middle slide index
        mid_idx = chunk[len(chunk)//2]["index"]
        for c in data:
            if not c.get("source_slide"):
                c["source_slide"] = mid_idx
        all_cards.extend(data)

        # gentle pacing to play nice with rate limits
        time.sleep(base_sleep)

    # Load optimization config for enhanced deduplication
    try:
        with open("config/flashcard_optimization_config.json", 'r') as f:
            optimization_config = json.load(f)
        dedup_config = optimization_config.get("duplicate_detection", {})
    except FileNotFoundError:
        dedup_config = {
            "fuzzy_match_threshold": 0.3,
            "semantic_similarity_threshold": 0.85,
            "context_weight": 0.2,
            "confidence_weight": 0.8
        }
    
    return dedupe_and_filter(all_cards, min_def_len=min_def_len, config=dedup_config)


# ---------------------------
# Output
# ---------------------------
def write_quizlet_tsv(cards: List[Dict[str, Any]], out_path: Path):
    out_path.parent.mkdir(parents=True, exist_ok=True)
    with out_path.open("w", encoding="utf-8", newline="") as f:
        w = csv.writer(f, delimiter="\t")
        for c in cards:
            w.writerow([c["term"], c["definition"]])


def write_json_log(cards: List[Dict[str, Any]], out_path: Path):
    with out_path.open("w", encoding="utf-8") as f:
        json.dump(cards, f, ensure_ascii=False, indent=2)


# ---------------------------
# Budget estimation (very rough)
# ---------------------------
def estimate_usd_for_chunked_calls(num_chunks: int, model: str, approx_input_toks=1200, approx_output_toks=350):
    """
    Very rough estimate based on per-window context size.
    Defaults assume window=3 with titles/bullets/notes.
    """
    pricing = {
        # input_per_1k, output_per_1k
        "gpt-4o-mini": (0.000150, 0.000600),
        "gpt-4o": (0.0025, 0.0100),
        "gpt-4.1-mini": (0.00030, 0.0012),
    }
    inp, outp = pricing.get(model, pricing["gpt-4o-mini"])
    per_call = (approx_input_toks / 1000.0) * inp + (approx_output_toks / 1000.0) * outp
    return per_call * num_chunks


# ---------------------------
# Main
# ---------------------------
def main():
    ap = argparse.ArgumentParser(description="Extract term–definition pairs from PPTX and output Quizlet-ready TSVs.")
    ap.add_argument("path", help="Folder containing .pptx files or a single .pptx file")
    ap.add_argument("--model", default=DEFAULT_MODEL)
    ap.add_argument("--window", type=int, default=3)
    ap.add_argument("--max-retries", type=int, default=4)
    ap.add_argument("--sleep", type=float, default=0.3, help="Base sleep between API calls")
    ap.add_argument("--min-def-len", type=int, default=MIN_DEF_LEN)
    ap.add_argument("--dry-run", action="store_true")
    ap.add_argument("--budget-usd", type=float, default=0.0, help="If >0, stop when rough estimate exceeds this.")
    ap.add_argument("--verbose", action="store_true")
    args = ap.parse_args()

    p = Path(args.path)
    if not p.exists():
        print(f"Path not found: {p}", file=sys.stderr)
        sys.exit(1)

    # Collect files
    files: List[Path] = []
    if p.is_dir():
        files = sorted([x for x in p.iterdir() if x.suffix.lower() == ".pptx"])
    else:
        if p.suffix.lower() != ".pptx":
            print("Provide a .pptx file or a folder containing .pptx files.", file=sys.stderr)
            sys.exit(1)
        files = [p]

    if not args.dry_run:
        # Init client (needs OPENAI_API_KEY)
        client = OpenAI()
    else:
        client = None  # type: ignore

    # Load cross-lesson data for context enhancement
    cross_lesson_data = None
    if not args.dry_run:
        cross_lesson_data = load_cross_lesson_data()
        if args.verbose and cross_lesson_data.get("content_index"):
            print(f"[context] Loaded cross-lesson data for {len(cross_lesson_data['content_index'])} lessons")

    total_estimate = 0.0
    for fp in files:
        if args.verbose:
            print(f"\n=== Processing: {fp.name} ===")

        # Parse once to count chunks for budget estimation
        prs = Presentation(str(fp))
        slides_ctx = []
        for idx, slide in enumerate(prs.slides, start=1):
            info = extract_slide_text(slide)
            info["index"] = idx
            slides_ctx.append(info)
        filtered = []
        for s in slides_ctx:
            if is_skippable_title(s["title"]) and not s["body"] and not s["notes"]:
                continue
            filtered.append(s)
        chunks = slide_windows(filtered, window=args.window)
        this_estimate = estimate_usd_for_chunked_calls(len(chunks), args.model)
        total_estimate += this_estimate

        if args.budget_usd > 0 and total_estimate > args.budget_usd:
            print(f"[budget] Estimated cost {total_estimate:.4f} exceeds budget {args.budget_usd:.4f}. Stopping.", file=sys.stderr)
            break

        if args.dry_run:
            if args.verbose:
                print(f"[dry-run] {len(chunks)} chunks; est ${this_estimate:.4f}")
            continue

        # Real extraction
        cards = extract_cards_from_presentation(
            fp,
            client=client,
            model=args.model,
            window=args.window,
            max_retries=args.max_retries,
            base_sleep=args.sleep,
            min_def_len=args.min_def_len,
            dry_run=False,
            verbose=args.verbose,
            cross_lesson_data=cross_lesson_data,
        )

        out_tsv = fp.with_suffix(".quizlet.tsv")
        out_json = fp.with_suffix(".quizlet.json")
        write_quizlet_tsv(cards, out_tsv)
        write_json_log(cards, out_json)

        # Also produce a cleaned copy (trims whitespace; stays TSV)
        # (Quizlet likes simple [term<TAB>definition] lines)
        # Already clean by construction, but we can rewrite to be safe:
        cleaned = []
        for c in cards:
            term = re.sub(r"\s+", " ", c["term"]).strip()
            defi = re.sub(r"\s+", " ", c["definition"]).strip()
            cleaned.append((term, defi))
        out_clean = fp.with_suffix(".quizlet.clean.tsv")
        with out_clean.open("w", encoding="utf-8", newline="") as f:
            w = csv.writer(f, delimiter="\t")
            for term, defi in cleaned:
                w.writerow([term, defi])

        if args.verbose:
            print(f"[done] {fp.name}: {len(cards)} cards →")
            print(f"       {out_tsv.name}, {out_clean.name}, {out_json.name}")

    if args.verbose:
        print(f"\nAll done. Estimated total spend: ${total_estimate:.4f}")


if __name__ == "__main__":
    main()
