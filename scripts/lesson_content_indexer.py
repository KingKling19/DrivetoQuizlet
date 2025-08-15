#!/usr/bin/env python3
"""
Lesson Content Indexer - Phase 1 Implementation
Extracts and indexes content from all lessons for cross-lesson context analysis.
"""

import os
import json
import pickle
import logging
from pathlib import Path
from typing import Dict, List, Tuple, Any
import numpy as np
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
import openai
from pptx import Presentation
import pandas as pd

# Setup logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class LessonContentIndexer:
    """Index and analyze lesson content for cross-lesson context system."""
    
    def __init__(self, workspace_path: str = "/workspace"):
        self.workspace_path = Path(workspace_path)
        self.lessons_path = self.workspace_path / "lessons"
        self.config_path = self.workspace_path / "config"
        self.config_path.mkdir(exist_ok=True)
        
        # Initialize storage
        self.content_index = {}
        self.tfidf_vectorizer = TfidfVectorizer(
            max_features=1000,
            stop_words='english',
            ngram_range=(1, 2)
        )
        self.semantic_embeddings = {}
        
        # Setup OpenAI
        self._setup_openai()
    
    def _setup_openai(self):
        """Setup OpenAI client for embeddings."""
        try:
            # Try to get API key from environment
            api_key = os.getenv('OPENAI_API_KEY')
            if api_key:
                openai.api_key = api_key
                logger.info("OpenAI API key loaded from environment")
            else:
                logger.warning("No OpenAI API key found - semantic embeddings will be skipped")
        except Exception as e:
            logger.error(f"Error setting up OpenAI: {e}")
    
    def extract_lesson_content(self, lesson_path: Path) -> Dict[str, Any]:
        """Extract all text content from a lesson directory."""
        content = {
            'lesson_id': lesson_path.name,
            'text_content': [],
            'key_concepts': [],
            'metadata': {
                'path': str(lesson_path),
                'files_processed': [],
                'content_sources': []
            }
        }
        
        # Process PowerPoint presentations
        presentations_path = lesson_path / "presentations"
        if presentations_path.exists():
            for ppt_file in presentations_path.glob("*.pptx"):
                try:
                    text = self._extract_ppt_text(ppt_file)
                    if text:
                        content['text_content'].append(text)
                        content['metadata']['files_processed'].append(str(ppt_file))
                        content['metadata']['content_sources'].append('presentation')
                except Exception as e:
                    logger.error(f"Error processing {ppt_file}: {e}")
        
        # Process notes
        notes_path = lesson_path / "notes"
        if notes_path.exists():
            for note_file in notes_path.glob("*.md"):
                try:
                    with open(note_file, 'r', encoding='utf-8') as f:
                        text = f.read()
                        content['text_content'].append(text)
                        content['metadata']['files_processed'].append(str(note_file))
                        content['metadata']['content_sources'].append('notes')
                except Exception as e:
                    logger.error(f"Error processing {note_file}: {e}")
        
        # Process README.md files in lesson root
        readme_file = lesson_path / "README.md"
        if readme_file.exists():
            try:
                with open(readme_file, 'r', encoding='utf-8') as f:
                    text = f.read()
                    content['text_content'].append(text)
                    content['metadata']['files_processed'].append(str(readme_file))
                    content['metadata']['content_sources'].append('readme')
            except Exception as e:
                logger.error(f"Error processing {readme_file}: {e}")
        
        # Process processed content
        processed_path = lesson_path / "processed"
        if processed_path.exists():
            for proc_file in processed_path.glob("*.json"):
                try:
                    with open(proc_file, 'r', encoding='utf-8') as f:
                        data = json.load(f)
                        if 'flashcards' in data:
                            # Extract text from flashcards
                            flashcard_text = []
                            for card in data['flashcards']:
                                if isinstance(card, dict):
                                    flashcard_text.append(f"Q: {card.get('question', '')} A: {card.get('answer', '')}")
                            if flashcard_text:
                                content['text_content'].append('\n'.join(flashcard_text))
                                content['metadata']['files_processed'].append(str(proc_file))
                                content['metadata']['content_sources'].append('flashcards')
                except Exception as e:
                    logger.error(f"Error processing {proc_file}: {e}")
        
        # Combine all text content
        combined_text = '\n'.join(content['text_content'])
        content['combined_text'] = combined_text
        
        # Extract key concepts (simple keyword extraction)
        content['key_concepts'] = self._extract_key_concepts(combined_text)
        
        return content
    
    def _extract_ppt_text(self, ppt_path: Path) -> str:
        """Extract text from PowerPoint presentation."""
        text_content = []
        try:
            prs = Presentation(ppt_path)
            for slide in prs.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text:
                        slide_text.append(shape.text)
                text_content.append('\n'.join(slide_text))
        except Exception as e:
            logger.error(f"Error extracting text from {ppt_path}: {e}")
        
        return '\n'.join(text_content)
    
    def _extract_key_concepts(self, text: str) -> List[str]:
        """Extract key concepts from text (simplified implementation)."""
        if not text:
            return []
        
        # Simple keyword extraction based on capitalized terms and common patterns
        import re
        
        # Find capitalized phrases (potential concepts)
        concepts = set()
        
        # Pattern for acronyms (2+ capital letters)
        acronyms = re.findall(r'\b[A-Z]{2,}\b', text)
        concepts.update(acronyms)
        
        # Pattern for title case phrases (up to 4 words)
        title_case = re.findall(r'\b[A-Z][a-z]+(?:\s+[A-Z][a-z]+){0,3}\b', text)
        concepts.update(title_case)
        
        # Filter out common words and short terms
        filtered_concepts = []
        common_words = {'The', 'This', 'That', 'With', 'From', 'When', 'Where', 'What', 'How', 'Why'}
        
        for concept in concepts:
            if len(concept) >= 3 and concept not in common_words:
                filtered_concepts.append(concept)
        
        return list(set(filtered_concepts))[:50]  # Limit to top 50 concepts
    
    def generate_content_fingerprint(self, content: Dict[str, Any]) -> np.ndarray:
        """Generate TF-IDF vector for content."""
        text = content['combined_text']
        if not text:
            return np.zeros(1000)  # Return zero vector if no content
        
        try:
            # Fit vectorizer if not already fitted
            if not hasattr(self.tfidf_vectorizer, 'vocabulary_'):
                # Get all lesson texts first to fit the vectorizer
                all_texts = []
                for lesson_dir in self.lessons_path.iterdir():
                    if lesson_dir.is_dir():
                        lesson_content = self.extract_lesson_content(lesson_dir)
                        if lesson_content['combined_text']:
                            all_texts.append(lesson_content['combined_text'])
                
                if all_texts:
                    self.tfidf_vectorizer.fit(all_texts)
                else:
                    return np.zeros(1000)
            
            # Transform the current text
            vector = self.tfidf_vectorizer.transform([text])
            return vector.toarray()[0]
        
        except Exception as e:
            logger.error(f"Error generating content fingerprint: {e}")
            return np.zeros(1000)
    
    def generate_semantic_embedding(self, content: Dict[str, Any]) -> List[float]:
        """Generate semantic embedding using OpenAI API."""
        text = content['combined_text']
        if not text or not openai.api_key:
            return []
        
        try:
            # Truncate text if too long (OpenAI has token limits)
            max_chars = 8000
            if len(text) > max_chars:
                text = text[:max_chars] + "..."
            
            response = openai.Embedding.create(
                model="text-embedding-ada-002",
                input=text
            )
            return response['data'][0]['embedding']
        
        except Exception as e:
            logger.error(f"Error generating semantic embedding: {e}")
            return []
    
    def index_all_lessons(self) -> Dict[str, Any]:
        """Index all lessons in the lessons directory."""
        logger.info("Starting lesson content indexing...")
        
        if not self.lessons_path.exists():
            logger.error(f"Lessons directory not found: {self.lessons_path}")
            return {}
        
        indexed_lessons = {}
        
        for lesson_dir in self.lessons_path.iterdir():
            if lesson_dir.is_dir():
                logger.info(f"Processing lesson: {lesson_dir.name}")
                
                # Extract content
                content = self.extract_lesson_content(lesson_dir)
                
                # Generate fingerprint
                content_fingerprint = self.generate_content_fingerprint(content)
                
                # Generate semantic embedding
                semantic_embedding = self.generate_semantic_embedding(content)
                
                # Store indexed content
                indexed_lessons[lesson_dir.name] = {
                    'lesson_id': lesson_dir.name,
                    'content_fingerprint': content_fingerprint.tolist() if isinstance(content_fingerprint, np.ndarray) else content_fingerprint,
                    'semantic_embedding': semantic_embedding,
                    'key_concepts': content['key_concepts'],
                    'content_metadata': content['metadata'],
                    'text_length': len(content['combined_text']),
                    'has_content': len(content['combined_text']) > 0
                }
        
        self.content_index = indexed_lessons
        logger.info(f"Indexed {len(indexed_lessons)} lessons")
        
        return indexed_lessons
    
    def save_index(self):
        """Save the content index to files."""
        # Save main index as JSON
        index_path = self.config_path / "lesson_content_index.json"
        with open(index_path, 'w', encoding='utf-8') as f:
            json.dump(self.content_index, f, indent=2)
        
        # Save semantic embeddings separately (pickle for numpy arrays)
        embeddings_path = self.config_path / "semantic_embeddings.pkl"
        embeddings_data = {
            lesson_id: data['semantic_embedding'] 
            for lesson_id, data in self.content_index.items()
            if data['semantic_embedding']
        }
        with open(embeddings_path, 'wb') as f:
            pickle.dump(embeddings_data, f)
        
        # Save TF-IDF vectorizer
        vectorizer_path = self.config_path / "tfidf_vectorizer.pkl"
        with open(vectorizer_path, 'wb') as f:
            pickle.dump(self.tfidf_vectorizer, f)
        
        logger.info(f"Index saved to {index_path}")
        logger.info(f"Embeddings saved to {embeddings_path}")
        logger.info(f"Vectorizer saved to {vectorizer_path}")
    
    def load_index(self) -> bool:
        """Load existing content index."""
        try:
            index_path = self.config_path / "lesson_content_index.json"
            if index_path.exists():
                with open(index_path, 'r', encoding='utf-8') as f:
                    self.content_index = json.load(f)
                logger.info(f"Loaded content index with {len(self.content_index)} lessons")
                return True
        except Exception as e:
            logger.error(f"Error loading index: {e}")
        
        return False


def main():
    """Main function to run the lesson content indexer."""
    indexer = LessonContentIndexer()
    
    # Try to load existing index
    if not indexer.load_index():
        logger.info("No existing index found, creating new index...")
        
        # Index all lessons
        indexed_lessons = indexer.index_all_lessons()
        
        if indexed_lessons:
            # Save the index
            indexer.save_index()
            
            # Print summary
            total_lessons = len(indexed_lessons)
            lessons_with_content = sum(1 for data in indexed_lessons.values() if data['has_content'])
            
            print(f"\nIndexing Complete!")
            print(f"Total lessons processed: {total_lessons}")
            print(f"Lessons with content: {lessons_with_content}")
            print(f"Lessons without content: {total_lessons - lessons_with_content}")
            
            # Show sample of key concepts
            all_concepts = []
            for data in indexed_lessons.values():
                all_concepts.extend(data['key_concepts'])
            
            unique_concepts = list(set(all_concepts))
            print(f"\nTotal unique concepts found: {len(unique_concepts)}")
            print(f"Sample concepts: {unique_concepts[:10]}")
        
        else:
            print("No lessons found to index")
    else:
        print(f"Loaded existing index with {len(indexer.content_index)} lessons")


if __name__ == "__main__":
    main()