#!/usr/bin/env python3
"""
Lesson Content Indexer

Extracts and indexes content from all lessons to enable cross-lesson context analysis.
Creates searchable content database with fingerprints for similarity detection.
"""

import argparse
import json
import os
import re
import sys
from collections import defaultdict
from pathlib import Path
from typing import Dict, List, Any, Optional, Tuple
import hashlib
import pickle
from datetime import datetime

try:
    from openai import OpenAI
except ImportError:
    print("ERROR: openai>=1.0.0 is required. pip install openai", file=sys.stderr)
    raise

try:
    from pptx import Presentation
except ImportError:
    print("ERROR: python-pptx is required. pip install python-pptx", file=sys.stderr)
    raise

try:
    import numpy as np
    from sklearn.feature_extraction.text import TfidfVectorizer
    from sklearn.metrics.pairwise import cosine_similarity
except ImportError:
    print("ERROR: scikit-learn is required. pip install scikit-learn", file=sys.stderr)
    raise


class LessonContentIndexer:
    """Indexes lesson content for cross-lesson context analysis."""
    
    def __init__(self, lessons_dir: Path = Path("lessons"), config_dir: Path = Path("config")):
        self.lessons_dir = Path(lessons_dir)
        self.config_dir = Path(config_dir)
        self.config_dir.mkdir(exist_ok=True)
        
        # Initialize OpenAI client
        self.openai_client = OpenAI()
        
        # Content storage
        self.content_index = {}
        self.tfidf_vectorizer = None
        self.tfidf_matrix = None
        self.semantic_embeddings = {}
        
        # Load existing index if available
        self.index_file = self.config_dir / "lesson_content_index.json"
        self.embeddings_file = self.config_dir / "semantic_embeddings.pkl"
        self.load_existing_index()
    
    def load_existing_index(self):
        """Load existing content index if available."""
        if self.index_file.exists():
            try:
                with open(self.index_file, 'r', encoding='utf-8') as f:
                    self.content_index = json.load(f)
                print(f"✓ Loaded existing index with {len(self.content_index)} lessons")
            except Exception as e:
                print(f"⚠️  Could not load existing index: {e}")
        
        if self.embeddings_file.exists():
            try:
                with open(self.embeddings_file, 'rb') as f:
                    self.semantic_embeddings = pickle.load(f)
                print(f"✓ Loaded existing semantic embeddings for {len(self.semantic_embeddings)} lessons")
            except Exception as e:
                print(f"⚠️  Could not load existing embeddings: {e}")
    
    def extract_lesson_content(self, lesson_dir: Path) -> Dict[str, Any]:
        """Extract all content from a lesson directory."""
        lesson_id = lesson_dir.name
        content = {
            "lesson_id": lesson_id,
            "lesson_name": lesson_dir.name.replace("_", " "),
            "content_sources": {},
            "key_concepts": [],
            "prerequisites": [],
            "related_lessons": [],
            "content_metadata": {
                "created": datetime.now().isoformat(),
                "last_updated": datetime.now().isoformat(),
                "content_types": []
            }
        }
        
        # Extract from presentations
        presentations_dir = lesson_dir / "presentations"
        if presentations_dir.exists():
            content["content_sources"]["presentations"] = self._extract_presentation_content(presentations_dir)
            content["content_metadata"]["content_types"].append("presentations")
        
        # Extract from notes
        notes_dir = lesson_dir / "notes"
        if notes_dir.exists():
            content["content_sources"]["notes"] = self._extract_notes_content(notes_dir)
            content["content_metadata"]["content_types"].append("notes")
        
        # Extract from processed content
        processed_dir = lesson_dir / "processed"
        if processed_dir.exists():
            content["content_sources"]["processed"] = self._extract_processed_content(processed_dir)
            content["content_metadata"]["content_types"].append("processed")
        
        # Extract from output content
        output_dir = lesson_dir / "output"
        if output_dir.exists():
            content["content_sources"]["output"] = self._extract_output_content(output_dir)
            content["content_metadata"]["content_types"].append("output")
        
        return content
    
    def _extract_presentation_content(self, presentations_dir: Path) -> Dict[str, Any]:
        """Extract content from PowerPoint presentations."""
        content = {"slides": [], "total_slides": 0, "text_content": ""}
        
        for pptx_file in presentations_dir.glob("*.pptx"):
            try:
                presentation = Presentation(pptx_file)
                slide_texts = []
                
                for i, slide in enumerate(presentation.slides):
                    slide_content = self._extract_slide_text(slide)
                    slide_texts.append(slide_content)
                    
                    content["slides"].append({
                        "slide_number": i + 1,
                        "title": slide_content.get("title", ""),
                        "body_text": slide_content.get("body_text", ""),
                        "tables": slide_content.get("tables", []),
                        "images": slide_content.get("images", [])
                    })
                
                content["total_slides"] += len(presentation.slides)
                content["text_content"] += "\n\n".join([
                    f"Slide {i+1}: {slide.get('title', '')}\n{slide.get('body_text', '')}"
                    for i, slide in enumerate(content["slides"])
                ])
                
            except Exception as e:
                print(f"⚠️  Error processing {pptx_file}: {e}")
        
        return content
    
    def _extract_slide_text(self, slide) -> Dict[str, Any]:
        """Extract text content from a PowerPoint slide."""
        content = {
            "title": "",
            "body_text": "",
            "tables": [],
            "images": []
        }
        
        # Extract title
        try:
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                content["title"] = slide.shapes.title.text or ""
        except Exception:
            pass
        
        # Extract body text and tables
        for shape in slide.shapes:
            try:
                if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                    text = shape.text or ""
                    if text and text != content["title"]:
                        content["body_text"] += text + "\n"
                
                if getattr(shape, "has_table", False) and shape.has_table:
                    table_data = []
                    for row in shape.table.rows:
                        row_data = []
                        for cell in row.cells:
                            row_data.append(cell.text or "")
                        table_data.append(row_data)
                    content["tables"].append(table_data)
                
                if hasattr(shape, "image") and shape.image:
                    content["images"].append(f"image_{len(content['images'])}")
                    
            except Exception:
                continue
        
        return content
    
    def _extract_notes_content(self, notes_dir: Path) -> Dict[str, Any]:
        """Extract content from notes (placeholder for future OCR implementation)."""
        content = {
            "note_files": [],
            "total_notes": 0,
            "text_content": ""
        }
        
        # For now, just collect file names
        # TODO: Implement OCR for handwritten notes
        for note_file in notes_dir.glob("*"):
            if note_file.suffix.lower() in ['.png', '.jpg', '.jpeg', '.pdf']:
                content["note_files"].append(note_file.name)
                content["total_notes"] += 1
        
        return content
    
    def _extract_processed_content(self, processed_dir: Path) -> Dict[str, Any]:
        """Extract content from processed files."""
        content = {
            "processed_files": [],
            "text_content": ""
        }
        
        for file_path in processed_dir.rglob("*"):
            if file_path.is_file() and file_path.suffix.lower() in ['.txt', '.json', '.md']:
                try:
                    with open(file_path, 'r', encoding='utf-8') as f:
                        file_content = f.read()
                        content["processed_files"].append({
                            "filename": file_path.name,
                            "content": file_content[:1000]  # Limit content length
                        })
                        content["text_content"] += file_content + "\n"
                except Exception as e:
                    print(f"⚠️  Error reading {file_path}: {e}")
        
        return content
    
    def _extract_output_content(self, output_dir: Path) -> Dict[str, Any]:
        """Extract content from output files (flashcards, etc.)."""
        content = {
            "output_files": [],
            "flashcards": [],
            "text_content": ""
        }
        
        for file_path in output_dir.rglob("*"):
            if file_path.is_file():
                if file_path.suffix.lower() == '.tsv':
                    # Extract flashcards from TSV files
                    try:
                        with open(file_path, 'r', encoding='utf-8') as f:
                            for line in f:
                                parts = line.strip().split('\t')
                                if len(parts) >= 2:
                                    content["flashcards"].append({
                                        "term": parts[0],
                                        "definition": parts[1]
                                    })
                                    content["text_content"] += f"{parts[0]}: {parts[1]}\n"
                    except Exception as e:
                        print(f"⚠️  Error reading {file_path}: {e}")
                
                content["output_files"].append(file_path.name)
        
        return content
    
    def generate_content_fingerprint(self, content: Dict[str, Any]) -> List[float]:
        """Generate TF-IDF fingerprint for lesson content."""
        # Combine all text content
        all_text = ""
        for source_type, source_content in content["content_sources"].items():
            if isinstance(source_content, dict) and "text_content" in source_content:
                all_text += source_content["text_content"] + "\n"
        
        # Clean and normalize text
        all_text = self._clean_text(all_text)
        
        # Generate TF-IDF vector
        if self.tfidf_vectorizer is None:
            self.tfidf_vectorizer = TfidfVectorizer(
                max_features=1000,
                stop_words='english',
                ngram_range=(1, 2),
                min_df=1
            )
            # Initialize with current text
            self.tfidf_vectorizer.fit([all_text])
        
        # Transform text to vector
        vector = self.tfidf_vectorizer.transform([all_text])
        return vector.toarray()[0].tolist()
    
    def generate_semantic_embedding(self, content: Dict[str, Any]) -> List[float]:
        """Generate semantic embedding using OpenAI API."""
        # Combine all text content
        all_text = ""
        for source_type, source_content in content["content_sources"].items():
            if isinstance(source_content, dict) and "text_content" in source_content:
                all_text += source_content["text_content"] + "\n"
        
        # Limit text length for API
        all_text = self._clean_text(all_text)[:8000]  # OpenAI limit
        
        try:
            response = self.openai_client.embeddings.create(
                model="text-embedding-3-small",
                input=all_text
            )
            return response.data[0].embedding
        except Exception as e:
            print(f"⚠️  Error generating semantic embedding: {e}")
            return []
    
    def _clean_text(self, text: str) -> str:
        """Clean and normalize text content."""
        # Remove extra whitespace
        text = re.sub(r'\s+', ' ', text)
        # Remove special characters but keep basic punctuation
        text = re.sub(r'[^\w\s\.\,\!\?\;\:\-\(\)]', '', text)
        return text.strip()
    
    def extract_key_concepts(self, content: Dict[str, Any]) -> List[str]:
        """Extract key concepts from lesson content."""
        concepts = set()
        
        # Extract from flashcards
        if "output" in content["content_sources"]:
            for flashcard in content["content_sources"]["output"].get("flashcards", []):
                concepts.add(flashcard["term"])
        
        # Extract from slide titles
        if "presentations" in content["content_sources"]:
            for slide in content["content_sources"]["presentations"].get("slides", []):
                title = slide.get("title", "")
                if title and len(title) > 3:
                    concepts.add(title)
        
        return list(concepts)[:50]  # Limit to top 50 concepts
    
    def index_all_lessons(self) -> Dict[str, Any]:
        """Index content from all lessons."""
        print("🔍 Indexing all lessons...")
        
        lesson_dirs = [d for d in self.lessons_dir.iterdir() if d.is_dir()]
        print(f"Found {len(lesson_dirs)} lesson directories")
        
        for lesson_dir in lesson_dirs:
            lesson_id = lesson_dir.name
            print(f"📚 Processing lesson: {lesson_id}")
            
            # Extract content
            content = self.extract_lesson_content(lesson_dir)
            
            # Generate fingerprints
            content["content_fingerprint"] = self.generate_content_fingerprint(content)
            content["semantic_embedding"] = self.generate_semantic_embedding(content)
            
            # Extract key concepts
            content["key_concepts"] = self.extract_key_concepts(content)
            
            # Store in index
            self.content_index[lesson_id] = content
            
            print(f"✓ Indexed {lesson_id} with {len(content['key_concepts'])} concepts")
        
        # Save index
        self.save_index()
        
        return self.content_index
    
    def save_index(self):
        """Save content index to files."""
        # Save main index
        with open(self.index_file, 'w', encoding='utf-8') as f:
            json.dump(self.content_index, f, indent=2, ensure_ascii=False)
        
        # Save semantic embeddings separately
        embeddings_data = {}
        for lesson_id, content in self.content_index.items():
            if "semantic_embedding" in content:
                embeddings_data[lesson_id] = content["semantic_embedding"]
        
        with open(self.embeddings_file, 'wb') as f:
            pickle.dump(embeddings_data, f)
        
        print(f"✓ Saved index with {len(self.content_index)} lessons")
    
    def get_lesson_content(self, lesson_id: str) -> Optional[Dict[str, Any]]:
        """Get content for a specific lesson."""
        return self.content_index.get(lesson_id)
    
    def get_all_lesson_ids(self) -> List[str]:
        """Get all lesson IDs in the index."""
        return list(self.content_index.keys())
    
    def search_content(self, query: str, top_k: int = 5) -> List[Tuple[str, float]]:
        """Search for lessons containing specific content."""
        if not self.content_index:
            return []
        
        # Simple text-based search for now
        results = []
        query_lower = query.lower()
        
        for lesson_id, content in self.content_index.items():
            score = 0
            text_content = ""
            
            # Combine all text content
            for source_type, source_content in content["content_sources"].items():
                if isinstance(source_content, dict) and "text_content" in source_content:
                    text_content += source_content["text_content"] + " "
            
            # Calculate relevance score
            if query_lower in text_content.lower():
                score += 1
            
            # Check key concepts
            for concept in content.get("key_concepts", []):
                if query_lower in concept.lower():
                    score += 2
            
            if score > 0:
                results.append((lesson_id, score))
        
        # Sort by score and return top k
        results.sort(key=lambda x: x[1], reverse=True)
        return results[:top_k]


def main():
    parser = argparse.ArgumentParser(description="Index lesson content for cross-lesson context analysis")
    parser.add_argument("--lessons-dir", default="lessons", help="Directory containing lessons")
    parser.add_argument("--config-dir", default="config", help="Directory for configuration files")
    parser.add_argument("--rebuild", action="store_true", help="Rebuild index from scratch")
    parser.add_argument("--search", help="Search for content in indexed lessons")
    
    args = parser.parse_args()
    
    # Initialize indexer
    indexer = LessonContentIndexer(
        lessons_dir=Path(args.lessons_dir),
        config_dir=Path(args.config_dir)
    )
    
    if args.search:
        # Search mode
        results = indexer.search_content(args.search)
        print(f"\n🔍 Search results for '{args.search}':")
        for lesson_id, score in results:
            print(f"  {lesson_id}: score {score}")
    else:
        # Index mode
        if args.rebuild:
            indexer.content_index = {}
            indexer.semantic_embeddings = {}
        
        # Index all lessons
        indexer.index_all_lessons()
        
        print(f"\n✅ Indexing complete!")
        print(f"📊 Indexed {len(indexer.content_index)} lessons")
        print(f"📁 Index saved to: {indexer.index_file}")
        print(f"🧠 Embeddings saved to: {indexer.embeddings_file}")


if __name__ == "__main__":
    main()
